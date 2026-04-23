"""
HIPAApotamus Pitch Deck Builder
--------------------------------
Design: "Ink & Brass" — editorial, minimal, legally-toned.

Click-1 sequence, per slide:
    Phase A (0–250ms)  the hero word glides up-left and fades while a small
                        dock version fades in at the top corner.  Motion uses
                        50% accel / 50% decel for an ease-in-out curve.
    Phase B (350–750ms) the first content block fades in (and rises slightly
                        on "rise" slides — subtle vertical easing).

Subsequent clicks fade the next content block in.  Each slide picks one of
several layouts (stack-left, stack-center, grid-2x2, two-col) so the eye
has a new pattern to land on each time.
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Design system
# ---------------------------------------------------------------------------

INK          = RGBColor(0x0B, 0x0E, 0x11)
BONE         = RGBColor(0xEC, 0xE6, 0xDA)
BONE_DIM     = RGBColor(0xC7, 0xC0, 0xB2)
SLATE        = RGBColor(0x8A, 0x8E, 0x96)
SLATE_DEEP   = RGBColor(0x5A, 0x5F, 0x67)
BRASS        = RGBColor(0xC9, 0x98, 0x56)
CRIMSON      = RGBColor(0xA8, 0x4A, 0x4A)
HAIR         = RGBColor(0x26, 0x2B, 0x31)

FONT_DISPLAY = "Georgia"
FONT_BODY    = "Helvetica"
FONT_MONO    = "Courier New"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Hero (center) and dock (top-left) positions — used to compute motion path.
HERO_BOX = (0.8, 2.5, 11.73, 2.0)         # (x, y, w, h) inches
HERO_FONT_PT = 88
DOCK_BOX = (0.8, 0.55, 9.0, 0.6)
DOCK_FONT_PT = 32

OUT_PATH = Path(__file__).parent / "HIPAApotamus_Pitch_Deck.pptx"

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ---------------------------------------------------------------------------
# Slide data
# ---------------------------------------------------------------------------
# Layout variants:
#   stack-left   — reveals stacked top-to-bottom, left-aligned
#   stack-center — reveals stacked top-to-bottom, centered
#   grid-2x2     — four reveals in a 2x2 grid (requires exactly 3 or 4 reveals)
#   two-col      — reveals split into two columns (left + right alternating)
# reveal_motion: "fade" (pure opacity) or "rise" (opacity + small y-translate)
# ---------------------------------------------------------------------------

SLIDES = [
    {
        "kind": "cover",
        "word": "HIPAApotamus",
        "tagline": "The HIPAA-native contract platform.",
        "sub": "Built for the Mississippi Department of Health.",
    },
    {
        "kind": "feature",
        "word": "Spreadsheets",
        "layout": "stack-left",
        "reveal_motion": "rise",
        "lede": "How most clinics track Business Associate Agreements today.",
        "reveals": [
            ["No attribution.",  "Who signed? From which IP? On what document version?",
             "45 CFR §164.312(b) — audit controls"],
            ["No integrity.",    "A row can be edited. A PDF on a shared drive can be replaced.",
             "45 CFR §164.312(c)(1) — protection from improper alteration"],
            ["No retention.",    "Six-year federal retention. Ten years under Mississippi law.",
             "45 CFR §164.530(j)(2)  ·  Miss. Code §41-9-69"],
            ["No audit packet.", "When OCR asks, the clinic gathers evidence by hand.",
             "— a weekend of panic, not a product —"],
        ],
    },
    {
        "kind": "feature",
        "word": "Emailing",
        "layout": "stack-left",
        "reveal_motion": "fade",
        "lede": "Delivering a BAA to a vendor looks like a solved problem — until you read the bill and the BAA.",
        "reveals": [
            ["DocuSign.",         "The industry default.",
             "", True],
            ["Expensive.",        "≈ $45 / user / month, per seat, plus per-envelope fees.",
             ""],
            ["Generic.",          "Doesn't know what a BAA is, or what §164.504(e) requires.",
             ""],
            ["Third-party.",      "Your PHI leaves the AWS BAA perimeter to reach their servers.",
             "— adds a Business Associate you must paper separately —"],
            ["Our answer.",       "Amazon SES inside the AWS BAA. Signed-URL invitations, no third party on the path.",
             "NIST SP 800-177  ·  AWS Healthcare Lens HCL_SEC7"],
        ],
    },
    {
        "kind": "feature",
        "word": "Verify",
        "layout": "grid-2x2",
        "reveal_motion": "rise",
        "lede": "Before a vendor sees the document, we prove they are the intended signer.",
        "reveals": [
            ["Magic link.",        "HMAC-SHA256 signed token, 72-hour TTL.",
             "src/lib/signing/token.ts"],
            ["Email OTP.",         "Six digits, ten-minute window, three-attempt lockout.",
             "src/lib/signing/otp.ts"],
            ["Authority to bind.", "OTP goes to the registered contact — not whoever forwarded the link.",
             "ESIGN §7001(c)  ·  UETA Miss. Code §75-12-1"],
            ["No vendor accounts.","Completion matters more than account creation friction.",
             ""],
        ],
    },
    {
        "kind": "feature",
        "word": "Signing",
        "layout": "grid-2x2",
        "reveal_motion": "fade",
        "lede": "The signing ceremony is four distinct, logged acts — not one button press.",
        "reveals": [
            ["Review.",  "Full contract text on screen; draft PDF downloadable before signing.",
             "ESIGN §7001(c)(1)(C)"],
            ["Consent.", "Explicit e-sign consent, timestamped before the canvas opens.",
             "ESIGN §7001(a)(1)  ·  UETA §75-12-5"],
            ["Capture.", "Drawn signature @ 2× DPR, IP, user-agent, consent time, method.",
             "src/components/signing/SigningInterface.tsx"],
            ["Seal.",    "SHA-256 of the final PDF, signed with AWS KMS.",
             "45 CFR §164.312(c)(1)-(2)"],
        ],
    },
    {
        "kind": "feature",
        "word": "Countersign",
        "layout": "stack-center",
        "reveal_motion": "rise",
        "lede": "A BAA is bilateral. The clinic signs too — or there is no contract.",
        "reveals": [
            ["Vendor first.",        "Business Associate signs via the OTP-gated link.",
             "status: pending_vendor_signature"],
            ["Clinic second.",       "HIPAA officer counter-signs from the dashboard.",
             "status: pending_countersignature"],
            ["Fully executed.",      "Both signatures embedded. Both hashes recorded. Contract is active.",
             "45 CFR §164.504(e) — required provisions satisfied"],
        ],
    },
    {
        "kind": "feature",
        "word": "Storage",
        "layout": "stack-left",
        "reveal_motion": "fade",
        "lede": "Signed documents belong in a place the clinic cannot accidentally delete.",
        "reveals": [
            ["Private S3.",          "Block-all-public-access. Server-side encryption. Versioning enabled.",
             "AWS Healthcare Lens HCL_SEC7"],
            ["Pre-signed access.",   "No public URLs — access is a short-lived signed URL, scoped to one object.",
             ""],
            ["Object Lock ready.",   "Bucket prepared for WORM retention in production (Compliance mode).",
             "10-year retention  ·  Miss. Code §41-9-69"],
            ["TLS 1.2+ enforced.",   "Bucket policy refuses non-SSL; SDK refuses plaintext.",
             "45 CFR §164.312(e)(1)"],
        ],
    },
    {
        "kind": "feature",
        "word": "Integrity",
        "layout": "grid-2x2",
        "reveal_motion": "rise",
        "lede": "How we answer the auditor's question: \"Has this document been altered since signing?\"",
        "reveals": [
            ["SHA-256.",             "256-bit fingerprint of the final PDF.",
             "src/app/api/baas/[id]/sign/route.ts"],
            ["KMS ECDSA_SHA_256.",   "Hash signed with an ECC NIST P-256 key on a FIPS 140-2 L3 HSM.",
             "src/lib/signing/kms.ts"],
            ["Verify every read.",   "Re-hash on download. Mismatch triggers a security alert.",
             "45 CFR §164.312(c)(1)"],
            ["Non-repudiation.",     "The signer cannot later claim \"that wasn't my document.\"",
             "— what makes it legally defensible —"],
        ],
    },
    {
        "kind": "feature",
        "word": "Tracking",
        "layout": "stack-left",
        "reveal_motion": "fade",
        "lede": "Every action in the signing ceremony is a first-class audit event.",
        "reveals": [
            ["Magic link accessed.",
             "Who opened the invitation, from where, when.",
             ""],
            ["OTP sent · verified · failed.",
             "Identity challenge recorded with attempts and lockouts.",
             ""],
            ["Consent · signature · hash · KMS seal.",
             "The four-act ceremony is reconstructable, frame by frame.",
             ""],
            ["Counter-signature · confirmation sent.",
             "Contract lifecycle is a log, not a guess.",
             "45 CFR §164.312(b) — audit controls"],
        ],
    },
    {
        "kind": "feature",
        "word": "Audit",
        "layout": "stack-left",
        "reveal_motion": "rise",
        "lede": "One ledger, surfaced three ways — for the operator, the auditor, and the regulator.",
        "reveals": [
            ["Agency view.",
             "Every audit event across every vendor in one filterable table — action, vendor, timestamp, performer. Filter by date range, action type, or free-text search.",
             "src/app/dashboard/history/page.tsx  ·  getRecentAuditLogs()"],
            ["Vendor drill-down.",
             "Click any vendor and the full chronological history opens — invitation, OTP, consent, signature, counter-sign, every reminder that ever went out.",
             "src/components/dashboard/BAADetailsModal.tsx  ·  getAuditLogsByVendor()"],
            ["Audit packet.",
             "One-click ZIP export: Executive Summary + signed BAA + Audit Trail PDF + Compliance Matrix. Hand it to OCR, the state health department, or a prospect's security review.",
             "src/lib/pdf/generator.ts  ·  AuditTrailPDF.tsx  ·  §164.312(b) evidence"],
        ],
    },
    {
        "kind": "feature",
        "word": "Reminders",
        "layout": "email-row",
        "reveal_motion": "fade",
        "lede": "Every status change sends a purpose-built email — and every email stays inside the AWS BAA perimeter.",
        "reveals": [
            [
                "Verification code.",
                "10-minute OTP sent to the signer's registered email.",
                "ESIGN §7001(c) — attribution",
                False, False,
                "assets/emails/verification_code.png",
            ],
            [
                "Signature reminder.",
                "Automatic nudges at 7 / 14 / 29 days since invitation.",
                "src/lib/reminders/policy.ts",
                False, False,
                "assets/emails/signature_reminder.png",
            ],
            [
                "Executed confirmation.",
                "Delivered to both parties when the clinic counter-signs.",
                "45 CFR §164.504(e) — fully executed",
                False, False,
                "assets/emails/executed_confirmation.png",
            ],
            [
                "Expiration notice.",
                "30 / 7 / 0-day warnings with a one-click renewal link.",
                "45 CFR §164.502(e)(1)(i)",
                False, False,
                "assets/emails/expiration_notice.png",
            ],
        ],
    },
    {
        "kind": "feature",
        "word": "Matrix",
        "layout": "two-col",
        "reveal_motion": "fade",
        "lede": "Every BAA clause is mapped to the regulation it satisfies — on the record, on the page.",
        "reveals": [
            ["45 CFR §164.504(e)(2).",
             "The required-elements checklist, mapped section by section.",
             "src/lib/baa/cfrMappings.ts"],
            ["MSDH sections.",
             "Mississippi-specific provisions cross-walked against state requirements.",
             "Compliance Matrix PDF — generated per BAA"],
            ["Evidence on demand.",
             "One export — ready for OCR, the state health department, or a prospect's security review.",
             "src/lib/pdf/ComplianceMatrixPDF.tsx"],
        ],
    },
    {
        "kind": "feature",
        "word": "Copilot",
        "layout": "stack-center",
        "reveal_motion": "rise",
        "lede": "The audit history is searchable in plain English — grounded in the real ledger.",
        "reveals": [
            ["Ask a question.",
             "\"Which BAAs expire in 30 days?\"   \"Show me ScriptGuard's signing history.\"",
             ""],
            ["Gemini 2.0 Flash.",
             "Streaming responses, grounded in audit logs — not model hallucination.",
             "src/lib/ai/gemini.ts"],
            ["RBAC enforced.",
             "The model sees only what the signed-in user is permitted to see.",
             "45 CFR §164.308(a)(4)"],
        ],
    },
    {
        "kind": "feature",
        "word": "Architecture",
        "layout": "stack-left",
        "reveal_motion": "fade",
        "lede": "A narrow, explainable AWS stack — every layer inside the BAA perimeter.",
        "reveals": [
            ["Identity.",
             "Amazon Cognito — MFA-ready, custom-attribute RBAC.",
             "cdk/lib/hipaa-baa-stack.ts"],
            ["Data.",
             "DynamoDB (PITR + streams)   ·   S3 (versioned, encrypted).",
             "AWS Healthcare Lens HCL_SEC1 / SEC7"],
            ["Cryptography.",
             "AWS KMS — FIPS 140-2 Level 3 HSM, ECC NIST P-256.",
             ""],
            ["Delivery.",
             "Amazon SES — email stays inside the AWS BAA perimeter.",
             ""],
            ["Intelligence.",
             "Google Gemini — grounded, streamed, scoped.",
             ""],
        ],
    },
    {
        "kind": "close",
        "word": "HIPAApotamus",
        "layout": "stack-center",
        "reveal_motion": "rise",
        "lede": "One HIPAA-native platform — not four SaaS tools duct-taped together.",
        "reveals": [
            ["DocuSign + Dropbox + Spreadsheets + Outlook.", "",
             "", True],
            ["One perimeter. One audit log. One matrix. One renewal loop.", "",
             ""],
            ["Built for Mississippi DOH. Extensible to every covered entity.", "",
             ""],
        ],
    },
]

# ---------------------------------------------------------------------------
# XML helpers
# ---------------------------------------------------------------------------

def _el(tag: str, **attrs) -> etree._Element:
    prefix, local = tag.split(":")
    ns = NS_P if prefix == "p" else NS_A
    e = etree.Element(etree.QName(ns, local), nsmap={"p": NS_P, "a": NS_A})
    for k, v in attrs.items():
        e.set(k, str(v))
    return e


def _sub(parent: etree._Element, tag: str, **attrs) -> etree._Element:
    prefix, local = tag.split(":")
    ns = NS_P if prefix == "p" else NS_A
    e = etree.SubElement(parent, etree.QName(ns, local))
    for k, v in attrs.items():
        e.set(k, str(v))
    return e


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color: RGBColor):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text: str, *,
             font=FONT_BODY, size=14, bold=False, italic=False,
             color=BONE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             spacing: Optional[float] = None):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n") if text else [""]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tx


def add_line(slide, x, y, w, color: RGBColor, weight_pt: float = 0.5):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, w, Emu(int(weight_pt * 12700))
    )
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------------------
# Chrome (always visible)
# ---------------------------------------------------------------------------

def build_chrome(slide, index: int, total: int, is_cover: bool = False) -> None:
    add_line(slide, Inches(0.8), Inches(7.02), Inches(11.73), HAIR, 0.5)
    if not is_cover:
        add_text(
            slide, Inches(0.8), Inches(7.12), Inches(6.0), Inches(0.25),
            "HIPAApotamus   ·   Feature walkthrough",
            font=FONT_BODY, size=9, color=SLATE_DEEP,
        )
        add_text(
            slide, Inches(7.0), Inches(7.12), Inches(5.53), Inches(0.25),
            f"{index:02d} / {total:02d}",
            font=FONT_MONO, size=9, color=SLATE_DEEP, align=PP_ALIGN.RIGHT,
        )


# ---------------------------------------------------------------------------
# Hero + dock (common to every non-cover slide)
# ---------------------------------------------------------------------------

def build_hero_and_dock(slide, data) -> Tuple[int, list, list]:
    """
    Returns (hero_word_id, hero_secondary_ids, dock_ids).

    The hero WORD is a single shape that physically moves + scales into the
    top-left position when click 1 fires.  There is NO separate dock-word
    shape — the hero itself becomes the small top-left title after motion.

    "Secondary" hero shapes are the centered brass rule and the centered
    italic lede, which fade out as the hero travels.  "Dock" shapes are the
    new top-left brass rule and top-left italic lede, which fade in near the
    end of the motion.
    """
    hx, hy, hw, hh = HERO_BOX
    hero = add_text(
        slide, Inches(hx), Inches(hy), Inches(hw), Inches(hh),
        data["word"],
        font=FONT_DISPLAY, size=HERO_FONT_PT, color=BONE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    hero_rule = add_rect(
        slide, Inches(6.17), Inches(hy + hh + 0.05), Inches(1.0),
        Emu(int(1.5 * 12700)), BRASS,
    )
    hero_secondary = [hero_rule.shape_id]
    if data.get("lede"):
        hero_lede = add_text(
            slide, Inches(0.8), Inches(hy + hh + 0.32), Inches(11.73), Inches(0.6),
            data["lede"],
            font=FONT_BODY, size=15, italic=True, color=SLATE, align=PP_ALIGN.CENTER,
        )
        hero_secondary.append(hero_lede.shape_id)

    # Dock decorations (hidden at start) — the hero word itself supplies the
    # word text at this position after it finishes motion + scale.
    dx, dy, dw, dh = DOCK_BOX
    # Brass rule sits slightly below where the scaled hero text lands
    dock_rule = add_rect(
        slide, Inches(dx + 0.02), Inches(1.30),
        Inches(0.55), Emu(int(1.5 * 12700)), BRASS,
    )
    dock_ids = [dock_rule.shape_id]
    if data.get("lede"):
        dock_lede = add_text(
            slide, Inches(dx), Inches(1.48), Inches(11.73), Inches(0.45),
            data["lede"],
            font=FONT_BODY, size=12, italic=True, color=SLATE, align=PP_ALIGN.LEFT,
        )
        dock_ids.append(dock_lede.shape_id)

    return hero.shape_id, hero_secondary, dock_ids


# ---------------------------------------------------------------------------
# Reveal builders — one per layout variant
# ---------------------------------------------------------------------------

def _reveal_block(slide, x, y, w, *, headline, subtext, citation,
                  strike=False, align=PP_ALIGN.LEFT, size_headline=22,
                  size_sub=13) -> List[int]:
    """Emit one reveal block at (x, y) with width w. Returns list of shape ids."""
    ids: List[int] = []
    if headline:
        th = add_text(
            slide, Inches(x), Inches(y), Inches(w), Inches(0.45),
            headline,
            font=FONT_BODY, size=size_headline, bold=True, color=BONE,
            align=align,
        )
        ids.append(th.shape_id)
        if strike:
            # Approximate strike width based on text length and size.
            char_w = 0.13 * (size_headline / 22.0)
            sw = min(char_w * len(headline), w - 0.1)
            # Align strike under the text run for the current alignment.
            if align == PP_ALIGN.CENTER:
                sx = x + (w - sw) / 2
            elif align == PP_ALIGN.RIGHT:
                sx = x + w - sw
            else:
                sx = x
            ln = add_line(
                slide, Inches(sx), Inches(y + 0.23), Inches(sw), CRIMSON, 1.6
            )
            ids.append(ln.shape_id)
    if subtext:
        ts = add_text(
            slide, Inches(x), Inches(y + 0.5), Inches(w), Inches(0.6),
            subtext,
            font=FONT_BODY, size=size_sub, color=BONE_DIM, align=align,
            spacing=1.25,
        )
        ids.append(ts.shape_id)
    if citation:
        tc = add_text(
            slide, Inches(x), Inches(y + 1.02), Inches(w), Inches(0.3),
            citation,
            font=FONT_MONO, size=10, italic=True, color=BRASS, align=align,
        )
        ids.append(tc.shape_id)
    return ids


def build_reveals_stack_left(slide, reveals) -> List[List[int]]:
    n = len(reveals)
    # Start lower to give dock+lede breathing room
    top = 1.95
    row_h = {1: 2.5, 2: 2.2, 3: 1.55, 4: 1.22, 5: 1.02}[min(n, 5)]
    groups: List[List[int]] = []
    for i, r in enumerate(reveals):
        groups.append(
            _reveal_block(
                slide,
                x=0.82, y=top + i * row_h, w=11.7,
                headline=r[0] if len(r) > 0 else "",
                subtext=r[1] if len(r) > 1 else "",
                citation=r[2] if len(r) > 2 else "",
                strike=r[3] if len(r) > 3 else False,
                align=PP_ALIGN.LEFT,
            )
        )
    return groups


def build_reveals_stack_center(slide, reveals) -> List[List[int]]:
    n = len(reveals)
    top = 2.2
    row_h = {1: 2.5, 2: 2.0, 3: 1.55, 4: 1.22, 5: 1.0}[min(n, 5)]
    groups: List[List[int]] = []
    # Narrower centered column for comfortable reading
    col_x, col_w = 1.7, 9.9
    for i, r in enumerate(reveals):
        groups.append(
            _reveal_block(
                slide,
                x=col_x, y=top + i * row_h, w=col_w,
                headline=r[0] if len(r) > 0 else "",
                subtext=r[1] if len(r) > 1 else "",
                citation=r[2] if len(r) > 2 else "",
                strike=r[3] if len(r) > 3 else False,
                align=PP_ALIGN.CENTER,
            )
        )
    return groups


def build_reveals_grid_2x2(slide, reveals) -> List[List[int]]:
    """4 reveals in 2x2. If only 3, third spans the bottom row."""
    n = len(reveals)
    assert n in (3, 4), "grid-2x2 expects 3 or 4 reveals"
    top = 2.1
    row_h = 2.35
    col_w = 5.55
    gutter = 0.65
    left_x = 0.82
    right_x = left_x + col_w + gutter
    groups: List[List[int]] = []
    if n == 4:
        coords = [
            (left_x,  top),
            (right_x, top),
            (left_x,  top + row_h),
            (right_x, top + row_h),
        ]
    else:
        coords = [
            (left_x,  top),
            (right_x, top),
            (left_x,  top + row_h),
        ]
    for (cx, cy), r in zip(coords, reveals):
        groups.append(
            _reveal_block(
                slide,
                x=cx, y=cy, w=col_w,
                headline=r[0] if len(r) > 0 else "",
                subtext=r[1] if len(r) > 1 else "",
                citation=r[2] if len(r) > 2 else "",
                strike=r[3] if len(r) > 3 else False,
                align=PP_ALIGN.LEFT,
                size_headline=21,
                size_sub=12,
            )
        )
    return groups


def build_reveals_two_col(slide, reveals) -> List[List[int]]:
    """Reveals alternate sides: 1st left, 2nd right, 3rd left, ..."""
    top = 2.0
    row_h = 1.75
    col_w = 5.4
    left_x = 0.82
    right_x = 13.333 - col_w - 0.82
    groups: List[List[int]] = []
    for i, r in enumerate(reveals):
        is_right = (i % 2 == 1)
        cx = right_x if is_right else left_x
        align = PP_ALIGN.RIGHT if is_right else PP_ALIGN.LEFT
        # Offset rows so left/right don't collide visually
        row_index = i // 2
        cy = top + row_index * row_h + (0 if not is_right else 0.35)
        groups.append(
            _reveal_block(
                slide,
                x=cx, y=cy, w=col_w,
                headline=r[0] if len(r) > 0 else "",
                subtext=r[1] if len(r) > 1 else "",
                citation=r[2] if len(r) > 2 else "",
                strike=r[3] if len(r) > 3 else False,
                align=align,
                size_headline=22,
                size_sub=13,
            )
        )
    return groups


def build_reveals_email_row(slide, reveals) -> List[List[int]]:
    """
    Horizontal row of up to 4 email thumbnails with captions below.

    Each reveal is expected to carry an image path at index 5:
        [headline, subtext, citation, strike, mono, image_path]

    On click N, the N-th thumbnail + its caption + its citation fade in
    together.  The hero word still animates to top-left on click 1.
    """
    n = len(reveals)
    assert 1 <= n <= 4, "email-row expects 1–4 reveals"

    slide_w_in = 13.333
    side_margin = 0.82
    gutter = 0.30
    usable = slide_w_in - 2 * side_margin
    col_w = (usable - gutter * (n - 1)) / n

    # Image block sized to feel screenshot-like without dominating the slide.
    img_h = 1.80
    row_top = 2.25

    groups: List[List[int]] = []
    here = Path(__file__).parent

    for i, r in enumerate(reveals):
        cx = side_margin + i * (col_w + gutter)
        headline  = r[0] if len(r) > 0 else ""
        subtext   = r[1] if len(r) > 1 else ""
        citation  = r[2] if len(r) > 2 else ""
        img_path  = r[5] if len(r) > 5 else ""

        shape_ids: List[int] = []

        # Image (the "screenshot")
        if img_path:
            full = here / img_path
            pic = slide.shapes.add_picture(
                str(full),
                Inches(cx), Inches(row_top),
                width=Inches(col_w), height=Inches(img_h),
            )
            shape_ids.append(pic.shape_id)

        # Caption block below the image
        caption_y = row_top + img_h + 0.22
        if headline:
            th = add_text(
                slide, Inches(cx), Inches(caption_y),
                Inches(col_w), Inches(0.35),
                headline,
                font=FONT_BODY, size=14, bold=True, color=BONE,
                align=PP_ALIGN.LEFT,
            )
            shape_ids.append(th.shape_id)
        if subtext:
            ts = add_text(
                slide, Inches(cx), Inches(caption_y + 0.38),
                Inches(col_w), Inches(0.9),
                subtext,
                font=FONT_BODY, size=11, color=BONE_DIM,
                align=PP_ALIGN.LEFT, spacing=1.2,
            )
            shape_ids.append(ts.shape_id)
        if citation:
            tc = add_text(
                slide, Inches(cx), Inches(caption_y + 1.25),
                Inches(col_w), Inches(0.3),
                citation,
                font=FONT_MONO, size=9, italic=True, color=BRASS,
                align=PP_ALIGN.LEFT,
            )
            shape_ids.append(tc.shape_id)

        groups.append(shape_ids)

    return groups


LAYOUT_BUILDERS = {
    "stack-left":   build_reveals_stack_left,
    "stack-center": build_reveals_stack_center,
    "grid-2x2":     build_reveals_grid_2x2,
    "two-col":      build_reveals_two_col,
    "email-row":    build_reveals_email_row,
}


# ---------------------------------------------------------------------------
# Cover / close
# ---------------------------------------------------------------------------

def build_cover(slide, data) -> dict:
    # Hero word, big, center
    hero = add_text(
        slide, Inches(0.8), Inches(3.0), Inches(11.73), Inches(1.3),
        data["word"],
        font=FONT_DISPLAY, size=74, color=BONE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    rule = add_rect(
        slide, Inches(6.27), Inches(4.18), Inches(0.8),
        Emu(int(1.5 * 12700)), BRASS,
    )
    tag = add_text(
        slide, Inches(0.8), Inches(4.5), Inches(11.73), Inches(0.5),
        data["tagline"],
        font=FONT_BODY, size=22, color=BONE_DIM, align=PP_ALIGN.CENTER,
    )
    sub = add_text(
        slide, Inches(0.8), Inches(5.05), Inches(11.73), Inches(0.4),
        data["sub"],
        font=FONT_BODY, size=13, italic=True, color=SLATE,
        align=PP_ALIGN.CENTER,
    )
    return {
        "kind": "cover",
        "hero_word_id": None,
        "hero_secondary_ids": [],
        "dock_ids": [],
        "click_groups": [[tag.shape_id, sub.shape_id]],
        "reveal_motion": "fade",
    }


def build_feature(slide, data) -> dict:
    layout = data.get("layout", "stack-left")
    hero_word_id, hero_secondary_ids, dock_ids = build_hero_and_dock(slide, data)
    builder = LAYOUT_BUILDERS[layout]
    click_groups = builder(slide, data.get("reveals", []))
    if not click_groups:
        click_groups = [[]]
    return {
        "kind": "feature",
        "hero_word_id": hero_word_id,
        "hero_secondary_ids": hero_secondary_ids,
        "dock_ids": dock_ids,
        "click_groups": click_groups,
        "reveal_motion": data.get("reveal_motion", "fade"),
    }


# ---------------------------------------------------------------------------
# Timing (animation) XML
# ---------------------------------------------------------------------------

class AnimIdCounter:
    def __init__(self, start: int = 2):
        self._n = start
    def next(self) -> int:
        self._n += 1
        return self._n


def _fade(parent, shape_id, ids, *, kind, node_type, delay_ms, duration):
    """Single fade entrance or exit effect."""
    par = _sub(parent, "p:par")
    ctn = _sub(
        par, "p:cTn",
        id=str(ids.next()),
        presetID="10", presetClass=kind, presetSubtype="0",
        fill="hold", grpId="0", nodeType=node_type,
    )
    _sub(_sub(ctn, "p:stCondLst"), "p:cond", delay=str(delay_ms))
    children = _sub(ctn, "p:childTnLst")

    # Snap visibility
    setel = _sub(children, "p:set")
    cbhvr = _sub(setel, "p:cBhvr")
    bctn = _sub(cbhvr, "p:cTn", id=str(ids.next()), dur="1", fill="hold")
    _sub(_sub(bctn, "p:stCondLst"), "p:cond", delay="0")
    _sub(_sub(cbhvr, "p:tgtEl"), "p:spTgt", spid=str(shape_id))
    anl = _sub(cbhvr, "p:attrNameLst")
    n = _sub(anl, "p:attrName"); n.text = "style.visibility"
    to = _sub(setel, "p:to")
    _sub(to, "p:strVal", val=("visible" if kind == "entr" else "hidden"))

    # Opacity tween
    anim = _sub(children, "p:anim", calcmode="lin", valueType="num")
    cbhvr2 = _sub(anim, "p:cBhvr", additive="base")
    _sub(cbhvr2, "p:cTn", id=str(ids.next()), dur=str(duration), fill="hold")
    _sub(_sub(cbhvr2, "p:tgtEl"), "p:spTgt", spid=str(shape_id))
    anl2 = _sub(cbhvr2, "p:attrNameLst")
    n2 = _sub(anl2, "p:attrName"); n2.text = "style.opacity"
    tav_lst = _sub(anim, "p:tavLst")
    if kind == "entr":
        t0, t1 = "0", "1"
    else:
        t0, t1 = "1", "0"
    tav0 = _sub(tav_lst, "p:tav", tm="0")
    _sub(_sub(tav0, "p:val"), "p:fltVal", val=t0)
    tav1 = _sub(tav_lst, "p:tav", tm="100000")
    _sub(_sub(tav1, "p:val"), "p:fltVal", val=t1)


def _motion(parent, shape_id, ids, *, dx, dy, node_type, delay_ms,
            duration, accel=50000, decel=50000):
    """Motion path effect with configurable easing (default: ease-in-out)."""
    par = _sub(parent, "p:par")
    ctn = _sub(
        par, "p:cTn",
        id=str(ids.next()),
        presetID="42", presetClass="path", presetSubtype="0",
        fill="hold", grpId="0", nodeType=node_type,
    )
    _sub(_sub(ctn, "p:stCondLst"), "p:cond", delay=str(delay_ms))
    children = _sub(ctn, "p:childTnLst")
    motion = _sub(
        children, "p:animMotion",
        origin="layout",
        path=f"M 0 0 L {dx:.4f} {dy:.4f} E",
        pathEditMode="relative", rAng="0", ptsTypes="",
    )
    cbhvr = _sub(motion, "p:cBhvr")
    inner = _sub(
        cbhvr, "p:cTn", id=str(ids.next()),
        dur=str(duration), fill="hold",
        accel=str(accel), decel=str(decel),
    )
    _sub(_sub(inner, "p:stCondLst"), "p:cond", delay="0")
    _sub(_sub(cbhvr, "p:tgtEl"), "p:spTgt", spid=str(shape_id))
    anl = _sub(cbhvr, "p:attrNameLst")
    n1 = _sub(anl, "p:attrName"); n1.text = "ppt_x"
    n2 = _sub(anl, "p:attrName"); n2.text = "ppt_y"


def _scale(parent, shape_id, ids, *, from_x=100000, from_y=100000,
           to_x, to_y, node_type, delay_ms, duration,
           accel=50000, decel=50000):
    """
    Grow/Shrink emphasis that scales the whole shape (text + decorations).
    Values are percent × 1000 (100000 = 100%, 34000 = 34%).
    Uses <p:from> + <p:to> so PowerPoint cannot misread the intent.
    """
    par = _sub(parent, "p:par")
    ctn = _sub(
        par, "p:cTn",
        id=str(ids.next()),
        presetID="6", presetClass="emph", presetSubtype="0",
        fill="hold", grpId="0", nodeType=node_type,
    )
    _sub(_sub(ctn, "p:stCondLst"), "p:cond", delay=str(delay_ms))
    children = _sub(ctn, "p:childTnLst")
    scale = _sub(children, "p:animScale")
    cbhvr = _sub(scale, "p:cBhvr")
    inner = _sub(
        cbhvr, "p:cTn", id=str(ids.next()),
        dur=str(duration), fill="hold",
        accel=str(accel), decel=str(decel),
    )
    _sub(_sub(inner, "p:stCondLst"), "p:cond", delay="0")
    _sub(_sub(cbhvr, "p:tgtEl"), "p:spTgt", spid=str(shape_id))
    _sub(scale, "p:from", x=str(from_x), y=str(from_y))
    _sub(scale, "p:to", x=str(to_x), y=str(to_y))


def _rise(parent, shape_id, ids, *, node_type, delay_ms, duration):
    """Small upward motion used on 'rise' reveals — runs alongside a fade."""
    _motion(parent, shape_id, ids,
            dx=0.0, dy=-0.02, node_type=node_type,
            delay_ms=delay_ms, duration=duration)


def build_timing(meta: dict) -> etree._Element:
    """
    Emit the <p:timing> tree for one slide.

    Click 1 (feature / close slide):
        Phase A (0 → 550ms) — the hero WORD itself physically travels and
        shrinks via animMotion + animScale.  At the same time, the centered
        brass rule and centered lede fade out (300ms), and the top-left
        brass rule and top-left lede fade in (150ms starting at 400ms).
        No cross-fade on the hero word — it's the same shape throughout.

        Phase B (700 → 1150ms) — the first content block fades in
        (with a small rise on 'rise' slides).

    Click 1 (cover): no hero transform, just fade the two taglines in.
    Click 2..N: fade the next content block in.
    """
    hero_word_id        = meta.get("hero_word_id")
    hero_secondary_ids  = meta.get("hero_secondary_ids", [])
    dock_ids            = meta.get("dock_ids", [])
    click_groups        = meta.get("click_groups", [])
    reveal_motion       = meta.get("reveal_motion", "fade")

    # Motion parameters for the hero word.  The target has its text centered
    # at roughly (1.6", 0.88") — matching where the top-left "dock" feels
    # like it belongs next to the brass rule at y = 1.3".  See HERO_BOX and
    # DOCK_BOX.
    MOTION_DX   = -0.38
    MOTION_DY   = -0.35
    # Final scale — 88pt → ~30pt is a 34% scale.
    SCALE_PCT   = 34000
    MOTION_DUR  = 550   # ms — long enough that the eye reads it as travel
    PHASE_B_DLY = 700   # ms — gap after motion before content reveals
    REVEAL_DUR  = 450   # ms

    ids = AnimIdCounter()

    timing = _el("p:timing")
    tn_lst = _sub(timing, "p:tnLst")
    root_par = _sub(tn_lst, "p:par")
    root_ctn = _sub(root_par, "p:cTn", id="1", dur="indefinite",
                    restart="never", nodeType="tmRoot")
    root_children = _sub(root_ctn, "p:childTnLst")

    seq = _sub(root_children, "p:seq", concurrent="1", nextAc="seek")
    seq_ctn = _sub(seq, "p:cTn", id="2", dur="indefinite", nodeType="mainSeq")
    main_children = _sub(seq_ctn, "p:childTnLst")

    def click_wrapper() -> etree._Element:
        outer_par = _sub(main_children, "p:par")
        outer_ctn = _sub(outer_par, "p:cTn", id=str(ids.next()), fill="hold")
        _sub(_sub(outer_ctn, "p:stCondLst"), "p:cond", delay="indefinite")
        outer_children = _sub(outer_ctn, "p:childTnLst")
        inner_par = _sub(outer_children, "p:par")
        inner_ctn = _sub(inner_par, "p:cTn", id=str(ids.next()), fill="hold")
        _sub(_sub(inner_ctn, "p:stCondLst"), "p:cond", delay="0")
        return _sub(inner_ctn, "p:childTnLst")

    def fade_group(container, shape_ids, *, delay_ms=0, duration=400,
                   node_start="withEffect", motion="fade"):
        first = True
        for sid in shape_ids:
            nt = node_start if first else "withEffect"
            _fade(container, sid, ids,
                  kind="entr", node_type=nt,
                  delay_ms=delay_ms, duration=duration)
            if motion == "rise":
                _rise(container, sid, ids,
                      node_type="withEffect",
                      delay_ms=delay_ms, duration=duration)
            first = False

    # --- Click 1 -----------------------------------------------------------
    if click_groups:
        container = click_wrapper()
        first_effect_used = False

        # Phase A — hero word travels + shrinks (same shape, no cross-fade)
        if hero_word_id is not None:
            _motion(container, hero_word_id, ids,
                    dx=MOTION_DX, dy=MOTION_DY,
                    node_type="clickEffect",
                    delay_ms=0, duration=MOTION_DUR)
            _scale(container, hero_word_id, ids,
                   from_x=100000, from_y=100000,
                   to_x=SCALE_PCT, to_y=SCALE_PCT,
                   node_type="withEffect",
                   delay_ms=0, duration=MOTION_DUR)
            first_effect_used = True

            # Hero decorations (centered rule + centered lede) fade out
            # during the first half of the motion so they don't clutter.
            for sid in hero_secondary_ids:
                _fade(container, sid, ids,
                      kind="exit", node_type="withEffect",
                      delay_ms=0, duration=300)

            # Dock decorations (top-left rule + top-left lede) fade in
            # just as the hero word is landing.
            dock_fade_delay = MOTION_DUR - 200   # start 200ms before landing
            for sid in dock_ids:
                _fade(container, sid, ids,
                      kind="entr", node_type="withEffect",
                      delay_ms=dock_fade_delay, duration=200)

        # Phase B — first content block reveals after the motion has settled
        first_group = click_groups[0]
        if first_group:
            nstart = "clickEffect" if not first_effect_used else "withEffect"
            fade_group(
                container, first_group,
                delay_ms=PHASE_B_DLY if first_effect_used else 0,
                duration=REVEAL_DUR,
                node_start=nstart,
                motion=reveal_motion,
            )

    # --- Click 2..N --------------------------------------------------------
    for grp in click_groups[1:]:
        if not grp:
            continue
        container = click_wrapper()
        fade_group(
            container, grp,
            delay_ms=0,
            duration=450,
            node_start="clickEffect",
            motion=reveal_motion,
        )

    # prev/next wiring
    prev_cl = _sub(seq, "p:prevCondLst")
    prev_cond = _sub(prev_cl, "p:cond", evt="onPrev", delay="0")
    _sub(_sub(prev_cond, "p:tgtEl"), "p:sldTgt")
    next_cl = _sub(seq, "p:nextCondLst")
    next_cond = _sub(next_cl, "p:cond", evt="onNext", delay="0")
    _sub(_sub(next_cond, "p:tgtEl"), "p:sldTgt")

    return timing


def attach_timing(slide, timing_el: etree._Element) -> None:
    sld = slide._element
    for e in sld.findall(f"{{{NS_P}}}timing"):
        sld.remove(e)
    sld.append(timing_el)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = len(SLIDES)

    for i, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        set_bg(slide, INK)
        is_cover = data["kind"] == "cover"
        build_chrome(slide, i, total, is_cover=is_cover)
        if is_cover:
            meta = build_cover(slide, data)
        else:
            meta = build_feature(slide, data)
        attach_timing(slide, build_timing(meta))

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
